import os
import io
import json
import zipfile
import xml.etree.ElementTree as ET
from typing import Dict, Any, Optional, List

def extract_txt_file(file_data: bytes, file_name: str) -> str:
    encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii']
    for encoding in encodings:
        try:
            return file_data.decode(encoding)
        except UnicodeDecodeError:
            continue
        except Exception as e:
            print(f"Unexpected error with encoding {encoding} for {file_name}: {e}")
            continue

    print(f"Warning: Could not decode {file_name} with standard encodings. Falling back to utf-8 with replacement.")
    return file_data.decode('utf-8', errors='replace')

def extract_pdf_file(file_data: bytes, file_name: str) -> str:
    try:
        import pdfplumber
        text = ""
        with pdfplumber.open(io.BytesIO(file_data)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        return text.strip()
    except Exception as e:
        print(f"PDFPlumber parsing failed for {file_name}: {e}")
        return f"Error: Could not extract text from PDF {file_name}"

def extract_docx_file(file_data: bytes, file_name: str) -> str:
    try:
        from unstructured.partition.docx import partition_docx
        elements = partition_docx(file=io.BytesIO(file_data))
        return "\n\n".join([str(el) for el in elements])
    except Exception as e:
        print(f"Unstructured DOCX parsing failed for {file_name}: {e}")

    try:
        from docx import Document
        doc = Document(io.BytesIO(file_data))
        text_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))

        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"python-docx parsing failed for {file_name}: {e}")
        return f"Error: Could not extract text from DOCX {file_name}"

def extract_doc_file(file_data: bytes, file_name: str) -> str:
    try:
        import docx2txt
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            text = docx2txt.process(tmp_file.name)
            os.unlink(tmp_file.name)  
        
        if text and text.strip():
            return text.strip()
        else:
            print(f"docx2txt returned empty text for {file_name}")
    except Exception as e:
        print(f"docx2txt parsing failed for {file_name}: {e}")
    
    try:
        import subprocess
        import tempfile
        
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            
            result = subprocess.run(['antiword', tmp_file.name], 
                                  capture_output=True, text=True, timeout=30)
            os.unlink(tmp_file.name)  
            
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
            else:
                print(f"antiword failed for {file_name}: {result.stderr}")
    except Exception as e:
        print(f"antiword parsing failed for {file_name}: {e}")
    
    try:
        from unstructured.partition.doc import partition_doc
        elements = partition_doc(file=io.BytesIO(file_data))
        text = "\n\n".join([str(el) for el in elements])
        if text and text.strip():
            return text.strip()
    except Exception as e:
        print(f"Unstructured DOC parsing failed for {file_name}: {e}")
    
    return f"Error: Could not extract text from DOC file {file_name}. Please convert to DOCX format for better compatibility."

def extract_pptx_file(file_data: bytes, file_name: str) -> str:
    try:
        from unstructured.partition.pptx import partition_pptx
        elements = partition_pptx(file=io.BytesIO(file_data))
        return "\n\n".join([str(el) for el in elements])
    except Exception as e:
        print(f"Unstructured PPTX parsing failed for {file_name}: {e}")

    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_data))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"--- Slide {slide_num} ---"
            text_parts.append(slide_text)

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"python-pptx parsing failed for {file_name}: {e}")
        return f"Error: Could not extract text from PPTX {file_name}"

def extract_ppt_file(file_data: bytes, file_name: str) -> str:
    try:
        from unstructured.partition.ppt import partition_ppt
        elements = partition_ppt(file=io.BytesIO(file_data))
        text = "\n\n".join([str(el) for el in elements])
        print("used unstructured")
        if text and text.strip():
            return text.strip()
    except Exception as e:
        print(f"Unstructured PPT parsing failed for {file_name}: {e}")
    
    try:
        import subprocess
        import tempfile
        import os
        
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            
            with tempfile.TemporaryDirectory() as tmp_dir:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', tmp_dir, tmp_file.name
                ], capture_output=True, text=True, timeout=60)
                
                if result.returncode == 0:
                    base_name = os.path.splitext(os.path.basename(tmp_file.name))[0]
                    txt_file = os.path.join(tmp_dir, f"{base_name}.txt")
                    
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8', errors='replace') as f:
                            text = f.read()
                        
                        os.unlink(tmp_file.name)  
                        
                        if text and text.strip():
                            return text.strip()
                
            os.unlink(tmp_file.name)  
            
    except Exception as e:
        print(f"LibreOffice conversion failed for {file_name}: {e}")
    
    return f"Error: Could not extract text from PPT file {file_name}. Please convert to PPTX format for better compatibility."


def extract_html_file(file_data: bytes, file_name: str) -> str:
    try:
        from unstructured.partition.html import partition_html
        elements = partition_html(file=io.BytesIO(file_data))
        return "\n\n".join([str(el) for el in elements])
    except Exception as e:
        print(f"Unstructured HTML parsing failed for {file_name}: {e}")

    try:
        from bs4 import BeautifulSoup
        text = extract_txt_file(file_data, file_name)
        soup = BeautifulSoup(text, 'html.parser')

        for script in soup(["script", "style"]):
            script.decompose()

        return soup.get_text(separator='\n\n', strip=True)
    except Exception as e:
        print(f"BeautifulSoup HTML parsing failed for {file_name}: {e}")
        return extract_txt_file(file_data, file_name)

def extract_rtf_file(file_data: bytes, file_name: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        rtf_content = extract_txt_file(file_data, file_name)
        return rtf_to_text(rtf_content)
    except Exception as e:
        print(f"RTF parsing failed for {file_name}: {e}")
        return extract_txt_file(file_data, file_name)

def extract_csv_file(file_data: bytes, file_name: str) -> str:
    try:
        import pandas as pd

        encodings = ['utf-8', 'latin-1', 'cp1252']
        df = None

        for encoding in encodings:
            try:
                df = pd.read_csv(io.BytesIO(file_data), encoding=encoding)
                break
            except UnicodeDecodeError:
                continue

        if df is None:
            df = pd.read_csv(io.BytesIO(file_data), encoding='utf-8', errors='replace')

        records = df.to_dict('records')

        text_parts = []
        text_parts.append(f"CSV File: {file_name}")
        text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
        text_parts.append(f"Total Rows: {len(df)}")
        text_parts.append("\nData Preview:")

        for i, record in enumerate(records[:10]): 
            text_parts.append(f"Row {i+1}: {json.dumps(record, default=str)}")

        if len(records) > 10:
            text_parts.append(f"... and {len(records) - 10} more rows")

        return "\n".join(text_parts)

    except Exception as e:
        print(f"CSV parsing failed for {file_name}: {e}")
        return extract_txt_file(file_data, file_name)

def extract_xlsx_file(file_data: bytes, file_name: str) -> str:
    try:
        import pandas as pd

        excel_file = pd.ExcelFile(io.BytesIO(file_data))
        text_parts = []
        text_parts.append(f"Excel File: {file_name}")
        text_parts.append(f"Sheets: {', '.join(excel_file.sheet_names)}")

        for sheet_name in excel_file.sheet_names:
            try:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)

                text_parts.append(f"\n--- Sheet: {sheet_name} ---")
                text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
                text_parts.append(f"Rows: {len(df)}")

                records = df.to_dict('records')
                text_parts.append("\nData Preview:")

                for i, record in enumerate(records[:5]):  
                    text_parts.append(f"Row {i+1}: {json.dumps(record, default=str)}")

                if len(records) > 5:
                    text_parts.append(f"... and {len(records) - 5} more rows")

            except Exception as e:
                text_parts.append(f"Error reading sheet {sheet_name}: {e}")

        return "\n".join(text_parts)

    except Exception as e:
        print(f"Excel parsing failed for {file_name}: {e}")
        return f"Error: Could not extract data from Excel file {file_name}"

def extract_opendocument_file(file_data: bytes, file_name: str) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(file_data), 'r') as zip_file:
            try:
                content_xml = zip_file.read('content.xml')
                text = extract_text_from_opendoc_xml(content_xml)
                
                if text.strip():
                    return text.strip()
                else:
                    print(f"No text content found in content.xml for {file_name}")
            except KeyError:
                print(f"content.xml not found in {file_name}")
            
            xml_files = ['styles.xml', 'meta.xml']
            extracted_texts = []
            
            for xml_file in xml_files:
                try:
                    xml_content = zip_file.read(xml_file)
                    text = extract_text_from_opendoc_xml(xml_content)
                    if text.strip():
                        extracted_texts.append(f"--- From {xml_file} ---\n{text}")
                except KeyError:
                    continue
                except Exception as e:
                    print(f"Error extracting from {xml_file}: {e}")
                    continue
            
            if extracted_texts:
                return "\n\n".join(extracted_texts)
                
    except zipfile.BadZipFile:
        print(f"Error: {file_name} is not a valid ZIP archive (corrupted OpenDocument file)")
    except Exception as e:
        print(f"Error extracting OpenDocument file {file_name}: {e}")
    
    return convert_opendoc_with_libreoffice(file_data, file_name)

def extract_text_from_opendoc_xml(xml_content: bytes) -> str:
    """
    Extract text from OpenDocument XML content.
    """
    try:
        root = ET.fromstring(xml_content)
        
        namespaces = {
            'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0',
            'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
            'draw': 'urn:oasis:names:tc:opendocument:xmlns:drawing:1.0',
            'presentation': 'urn:oasis:names:tc:opendocument:xmlns:presentation:1.0',
            'table': 'urn:oasis:names:tc:opendocument:xmlns:table:1.0'
        }
        
        text_parts = []
        
        for text_elem in root.iter():
            if text_elem.text and text_elem.text.strip():
                text_parts.append(text_elem.text.strip())
            if text_elem.tail and text_elem.tail.strip():
                text_parts.append(text_elem.tail.strip())
        
        try:
            for page in root.findall('.//draw:page', namespaces):
                page_texts = []
                for text_box in page.findall('.//draw:text-box', namespaces):
                    for p in text_box.findall('.//text:p', namespaces):
                        if p.text:
                            page_texts.append(p.text)
                if page_texts:
                    text_parts.extend(page_texts)
            
            for p in root.findall('.//text:p', namespaces):
                if p.text and p.text.strip():
                    text_parts.append(p.text.strip())
            
            for cell in root.findall('.//table:table-cell', namespaces):
                for p in cell.findall('.//text:p', namespaces):
                    if p.text and p.text.strip():
                        text_parts.append(p.text.strip())
                        
        except Exception as e:
            print(f"Error in structured OpenDocument extraction: {e}")
        
        unique_texts = []
        seen = set()
        for text in text_parts:
            if text not in seen and len(text) > 1: 
                unique_texts.append(text)
                seen.add(text)
        
        return "\n\n".join(unique_texts) if unique_texts else ""
        
    except ET.ParseError as e:
        print(f"XML parsing error: {e}")
        return ""
    except Exception as e:
        print(f"Error extracting text from XML: {e}")
        return ""

def convert_opendoc_with_libreoffice(file_data: bytes, file_name: str) -> str:
    try:
        import subprocess
        import tempfile
        
        file_ext = os.path.splitext(file_name.lower())[1]
        
        with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            
            with tempfile.TemporaryDirectory() as tmp_dir:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', tmp_dir, tmp_file.name
                ], capture_output=True, text=True, timeout=60)
                
                if result.returncode == 0:
                    base_name = os.path.splitext(os.path.basename(tmp_file.name))[0]
                    txt_file = os.path.join(tmp_dir, f"{base_name}.txt")
                    
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8', errors='replace') as f:
                            text = f.read()
                        
                        if text and text.strip():
                            return text.strip()
                else:
                    print(f"LibreOffice conversion failed: {result.stderr}")
                
        os.unlink(tmp_file.name)  
        
    except subprocess.TimeoutExpired:
        print(f"LibreOffice conversion timed out for {file_name}")
    except Exception as e:
        print(f"LibreOffice conversion failed for {file_name}: {e}")
    
    return f"Error: Could not extract text from OpenDocument file {file_name}. The file may be corrupted or contain only images/graphics."

def extract_unstructured_enhanced(file_data: bytes, file_name: str) -> str:
    file_extension = os.path.splitext(file_name.lower())[1]
    
    if file_extension in ['.odt', '.odp', '.ods']:
        return extract_opendocument_file(file_data, file_name)
    
    try:
        from unstructured.partition.auto import partition
        elements = partition(file=io.BytesIO(file_data), file_filename=file_name)
        text = "\n\n".join([str(el) for el in elements])
        
        if text.strip():
            return text.strip()
        else:
            print(f"Unstructured returned empty text for {file_name}")
            
    except Exception as e:
        print(f"Error during unstructured parsing for {file_name}: {e}")
    
    return extract_txt_file(file_data, file_name)

def extract_text_from_file(file_data: bytes, file_name: str) -> str:
    file_extension = os.path.splitext(file_name.lower())[1]
    print(f"Extracting text from '{file_name}' with extension '{file_extension}'")

    extractors = {
        '.txt': extract_txt_file,
        '.pdf': extract_pdf_file,
        '.docx': extract_docx_file,
        '.doc': extract_doc_file,
        '.pptx': extract_pptx_file,
        '.ppt': extract_ppt_file, 
        '.html': extract_html_file,
        '.htm': extract_html_file,
        '.rtf': extract_rtf_file,
        '.csv': extract_csv_file,
        '.xlsx': extract_xlsx_file,
        '.odt': extract_opendocument_file,
        '.odp': extract_opendocument_file,
        '.ods': extract_opendocument_file,
    }

    if file_extension in extractors:
        return extractors[file_extension](file_data, file_name)

    elif file_extension in ['.odt', '.odp', '.ods', '.md', '.rst', '.xml', '.json']:
        return extract_unstructured_enhanced(file_data, file_name)

    else:
        print(f"Unsupported file format '{file_extension}'. Attempting text extraction as a fallback.")
        return extract_txt_file(file_data, file_name)

def get_supported_file_types() -> List[str]:
    """
    Returns a list of supported file extensions.
    """
    return [
        '.txt', '.pdf', '.docx', '.doc', '.pptx', '.ppt',
        '.html', '.rtf', '.csv', '.xlsx',
        '.odt', '.odp', '.md', '.json'
    ]
